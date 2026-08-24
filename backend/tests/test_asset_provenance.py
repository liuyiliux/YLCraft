from __future__ import annotations

from pathlib import Path

from app.services.asset_provenance import clean_file, inspect_file


def test_text_provenance_preview_and_cleaning(tmp_path: Path):
    source = tmp_path / "note.txt"
    target = tmp_path / "note-cleaned.txt"
    source.write_text("标题\u200b：\u202e正文", encoding="utf-8")

    report = inspect_file(source, "text/plain")
    assert report.cleanable is True
    assert report.invisible_character_count == 1
    assert report.bidi_control_count == 1

    cleaned = clean_file(source, target)
    assert source.read_text(encoding="utf-8") == "标题\u200b：\u202e正文"
    assert target.read_text(encoding="utf-8") == "标题：正文"
    assert cleaned.cleanable is False


def test_image_metadata_is_removed_without_overwriting_source(tmp_path: Path):
    from PIL import Image, PngImagePlugin

    source = tmp_path / "image.png"
    target = tmp_path / "image-cleaned.png"
    info = PngImagePlugin.PngInfo()
    info.add_text("Comment", "generated metadata")
    Image.new("RGB", (4, 3), "white").save(source, pnginfo=info)

    report = inspect_file(source, "image/png")
    assert report.cleanable is True
    clean_file(source, target)

    assert source.exists() and target.exists()
    assert "Comment" in inspect_file(source, "image/png").metadata_keys
    assert "Comment" not in inspect_file(target, "image/png").metadata_keys


def test_unsupported_format_is_audit_only(tmp_path: Path):
    source = tmp_path / "clip.xyz"
    target = tmp_path / "clip-cleaned.xyz"
    source.write_bytes(b"opaque bytes")

    report = inspect_file(source, "application/x-unknown")
    assert report.supported is False
    cleaned = clean_file(source, target)
    assert target.read_bytes() == source.read_bytes()
    assert cleaned.cleanable is False


def test_video_audio_are_recognized_for_cleaning(tmp_path: Path):
    video = tmp_path / "clip.mp4"
    video.write_bytes(b"fake")
    audio = tmp_path / "track.mp3"
    audio.write_bytes(b"fake")

    video_report = inspect_file(video, "video/mp4")
    assert video_report.supported is True
    assert video_report.media_kind == "video"
    assert video_report.cleanable is True

    audio_report = inspect_file(audio, "audio/mpeg")
    assert audio_report.supported is True
    assert audio_report.media_kind == "audio"
    assert audio_report.cleanable is True


def test_pdf_metadata_is_removed_without_overwriting_source(tmp_path: Path):
    from pypdf import PdfReader, PdfWriter

    source = tmp_path / "doc.pdf"
    target = tmp_path / "doc-cleaned.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.add_metadata({"/Author": "SomeAuthor", "/Title": "MyDoc"})
    with open(source, "wb") as fh:
        writer.write(fh)

    report = inspect_file(source, "application/pdf")
    assert report.supported is True
    assert report.media_kind == "document"
    assert report.cleanable is True
    assert "author" in report.metadata_keys

    cleaned = clean_file(source, target)
    assert source.exists() and target.exists()
    assert "author" in inspect_file(source, "application/pdf").metadata_keys  # source untouched
    assert cleaned.cleanable is False
    assert "author" not in inspect_file(target, "application/pdf").metadata_keys
    # 内容与页数保留
    assert len(PdfReader(str(target)).pages) == 1


def test_docx_metadata_is_removed_without_overwriting_source(tmp_path: Path):
    import docx

    source = tmp_path / "doc.docx"
    target = tmp_path / "doc-cleaned.docx"
    document = docx.Document()
    document.add_paragraph("hello world")
    document.core_properties.author = "SomeAuthor"
    document.core_properties.title = "MyDoc"
    document.core_properties.keywords = "kw1;kw2"
    document.save(str(source))

    report = inspect_file(
        source,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )
    assert report.supported is True
    assert report.media_kind == "document"
    assert report.cleanable is True
    assert "author" in report.metadata_keys

    cleaned = clean_file(source, target)
    assert source.exists() and target.exists()
    assert "author" in inspect_file(source).metadata_keys  # source untouched
    assert cleaned.cleanable is False
    assert "author" not in inspect_file(target).metadata_keys
    # 内容保留
    assert docx.Document(str(target)).paragraphs[0].text == "hello world"


def test_pdf_docx_extensions_are_detected_by_suffix(tmp_path: Path):
    pdf = tmp_path / "a.pdf"
    pdf.write_bytes(b"%PDF-1.4")
    docx = tmp_path / "b.docx"
    docx.write_bytes(b"PK")

    assert inspect_file(pdf).media_kind == "document"
    assert inspect_file(docx).media_kind == "document"


def _write_zip_container(path: Path, payloads: dict[str, bytes]) -> None:
    """构造一个最小 zip 容器（docx/xlsx/pptx/odt/epub 共用）。"""
    import zipfile

    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for name, data in payloads.items():
            z.writestr(name, data)


def test_layer_a_unicode_scrub_is_comprehensive(tmp_path: Path):
    """Layer A 应覆盖远超早期实现的隐形/bidi/标签/非字符/空间同形字。"""
    # 组合各类隐形载体：零宽、bidi、软连字符、标签字符、非字符、保留可忽略符、私有用途、空间同形字
    source = tmp_path / "mixed.txt"
    target = tmp_path / "mixed-cleaned.txt"
    text = "a\u200b\u202e\u00ad\U000E0001\ufdd0\u2065\ue000\ue000\uff1c\u3000b"
    # 注意：U+FF1C 是全角小于号（空间/全角同形字），此处用私有用例验证
    source.write_text(text, encoding="utf-8")

    cleaned = clean_file(source, target)
    assert source.read_text(encoding="utf-8") == text  # 源文件不动
    assert cleaned.cleanable is False, cleaned.unicode_breakdown
    # 隐形载体应被清除，普通 ASCII 保留
    out = target.read_text(encoding="utf-8")
    assert "a" in out and "b" in out


def test_layer_a_preserves_emoji_glue_and_script_joiners(tmp_path: Path):
    """emoji ZWJ/VS 序列、以及合法连接脚本内的 ZWJ 不应被误删。"""
    from app.services.asset_provenance.service import _clean_text_layer_a

    cleaned, counts = _clean_text_layer_a("👨\u200d👩\u200d👧")
    assert "👨\u200d👩" in cleaned  # 家庭 emoji 序列保留
    assert counts == {}

    cleaned, counts = _clean_text_layer_a("❤\ufe0f")
    assert cleaned == "❤\ufe0f"

    # 波斯语正交 ZWJ 保留
    cleaned, counts = _clean_text_layer_a("می\u200cروم")
    assert "می\u200cروم" in cleaned


def test_layer_a_normalizes_space_homoglyphs(tmp_path: Path):
    """空间同形字应归一化为普通空格，并计入 breakdown。"""
    from app.services.asset_provenance.service import _clean_text_layer_a

    cleaned, counts = _clean_text_layer_a("a\u3000b\u00a0c")
    assert cleaned == "a b c"
    assert counts.get("space_homoglyph") == 2


def test_xlsx_metadata_is_removed(tmp_path: Path):
    import openpyxl

    source = tmp_path / "book.xlsx"
    target = tmp_path / "book-cleaned.xlsx"
    wb = openpyxl.Workbook()
    wb.active["A1"] = "hello"
    wb.properties.creator = "SomeAuthor"
    wb.properties.title = "MySheet"
    wb.save(str(source))

    report = inspect_file(
        source, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
    assert report.media_kind == "document"
    assert report.cleanable is True
    assert "author" in report.metadata_keys

    cleaned = clean_file(source, target)
    assert cleaned.cleanable is False, cleaned.metadata_keys
    assert "author" not in inspect_file(target).metadata_keys
    assert openpyxl.load_workbook(target).active["A1"].value == "hello"


def test_pptx_metadata_is_removed(tmp_path: Path):
    from pptx import Presentation

    source = tmp_path / "deck.pptx"
    target = tmp_path / "deck-cleaned.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.core_properties.author = "SomeAuthor"
    prs.save(str(source))

    report = inspect_file(
        source, "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert report.media_kind == "document"
    assert report.cleanable is True

    cleaned = clean_file(source, target)
    assert cleaned.cleanable is False, cleaned.metadata_keys
    assert "author" not in inspect_file(target).metadata_keys


def test_odt_metadata_is_removed(tmp_path: Path):
    import zipfile
    from xml.etree import ElementTree as ET

    source = tmp_path / "doc.odt"
    meta = (
        '<?xml version="1.0"?>'
        '<office:document-meta xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0" '
        'xmlns:dc="http://purl.org/dc/elements/1.1/" '
        'xmlns:meta="urn:oasis:names:tc:opendocument:xmlns:meta:1.0">'
        "<office:meta><dc:creator>SomeAuthor</dc:creator><dc:title>MyDoc</dc:title>"
        "<meta:generator>Writer</meta:generator></office:meta></office:document-meta>"
    )
    _write_zip_container(
        source,
        {
            "mimetype": b"application/vnd.oasis.opendocument.text",
            "content.xml": b'<?xml version="1.0"?><office:document-content xmlns:office="urn:oasis:names:tc:opendocument:xmlns:office:1.0"/>',
            "meta.xml": meta.encode(),
        },
    )

    report = inspect_file(source, "application/vnd.oasis.opendocument.text")
    assert report.media_kind == "document"
    assert report.cleanable is True
    assert "creator" in report.metadata_keys

    target = tmp_path / "doc-cleaned.odt"
    cleaned = clean_file(source, target)
    assert cleaned.cleanable is False, cleaned.metadata_keys
    with zipfile.ZipFile(target) as z:
        root = ET.fromstring(z.read("meta.xml"))
        creators = [
            e for e in root.iter()
            if "purl.org/dc" in e.tag and e.tag.split("}")[-1] == "creator"
        ]
        assert not creators


def test_epub_metadata_is_removed(tmp_path: Path):
    import zipfile
    from xml.etree import ElementTree as ET

    source = tmp_path / "book.epub"
    opf = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        '<package xmlns="http://www.idpf.org/2007/opf" '
        'xmlns:dc="http://purl.org/dc/elements/1.1/" version="3.0">'
        "<metadata><dc:title>MyBook</dc:title><dc:creator>SomeAuthor</dc:creator></metadata>"
        "<manifest/><spine/></package>"
    )
    container = (
        '<?xml version="1.0"?>'
        '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
        '<rootfiles><rootfile full-path="OEBPS/content.opf" '
        'media-type="application/oebps-package+xml"/></rootfiles></container>'
    )
    _write_zip_container(
        source,
        {"META-INF/container.xml": container.encode(), "OEBPS/content.opf": opf.encode()},
    )

    report = inspect_file(source, "application/epub+zip")
    assert report.media_kind == "document"
    assert report.cleanable is True
    assert "creator" in report.metadata_keys and "title" in report.metadata_keys

    target = tmp_path / "book-cleaned.epub"
    cleaned = clean_file(source, target)
    assert cleaned.cleanable is False, cleaned.metadata_keys
    with zipfile.ZipFile(target) as z:
        root = ET.fromstring(z.read("OEBPS/content.opf"))
        dc = [e for e in root.iter() if "purl.org/dc/elements" in e.tag]
        assert not dc


def test_gif_comment_metadata_is_removed(tmp_path: Path):
    from PIL import Image

    source = tmp_path / "anim.gif"
    target = tmp_path / "anim-cleaned.gif"
    im = Image.new("P", (4, 4))
    im.info["comment"] = b"gen"
    im.save(source, format="GIF")

    report = inspect_file(source, "image/gif")
    assert report.media_kind == "image"
    assert report.cleanable is True
    assert "comment" in report.metadata_keys

    cleaned = clean_file(source, target)
    assert cleaned.cleanable is False, cleaned.metadata_keys
    assert "comment" not in inspect_file(target).metadata_keys


def test_new_document_and_image_extensions_by_suffix(tmp_path: Path):
    for name in ["a.xlsx", "a.pptx", "a.odt", "a.epub", "a.gif", "a.avif"]:
        p = tmp_path / name
        p.write_bytes(b"PK" if name.endswith((".xlsx", ".pptx", ".epub")) else b"data")
        report = inspect_file(p)
        assert report.media_kind in {"document", "image"}, (name, report.media_kind)


def test_deep_watermark_detect_is_read_only_and_reports_ctrlregen(tmp_path: Path):
    from PIL import Image

    from app.services.asset_provenance import detect_deep_watermark

    source = tmp_path / "img.png"
    Image.new("RGB", (64, 64), "#4a7bb5").save(source)
    before = source.read_bytes()

    result = detect_deep_watermark(source, "image/png")
    assert result.supported is True
    assert result.media_kind == "image"
    assert "score" in result.ctrlregen
    assert "confidence" in result.ctrlregen
    assert result.ctrlregen["method"] == "statistical-ctrlregen-like"
    assert result.synthid["status"] == "skipped"
    # 只读：文件字节不变
    assert source.read_bytes() == before
    # 不伪装成已清理（任何 note 都不宣称能“移除”水印）
    assert all("移除" not in n for n in result.notes)
    # 明确只读上报
    assert any("未修改" in n for n in result.notes)


def test_deep_watermark_detect_image_with_noise_scores_higher(tmp_path: Path):
    """带噪/高频内容的图应得到更高（更强嵌入痕迹）得分，验证统计检测区分度。"""
    import random

    from PIL import Image

    from app.services.asset_provenance import detect_deep_watermark

    flat = tmp_path / "flat.png"
    noisy = tmp_path / "noisy.png"
    Image.new("RGB", (64, 64), (80, 80, 80)).save(flat)
    Image.new("RGB", (64, 64), (80, 80, 80)).save(noisy)
    img = Image.open(noisy)
    rnd = random.Random(42)
    for y in range(img.height):
        for x in range(img.width):
            img.putpixel((x, y), (rnd.randint(0, 255), rnd.randint(0, 255), rnd.randint(0, 255)))
    img.save(noisy)

    flat_result = detect_deep_watermark(flat, "image/png")
    noisy_result = detect_deep_watermark(noisy, "image/png")
    assert noisy_result.ctrlregen["score"] > flat_result.ctrlregen["score"]


def test_deep_watermark_detect_unsupported_and_synthid_skipped(tmp_path: Path):
    from app.services.asset_provenance import detect_deep_watermark

    # 文本：返回 unsupported 且不伪装
    text = tmp_path / "note.txt"
    text.write_text("hello", encoding="utf-8")
    result = detect_deep_watermark(text, "text/plain")
    assert result.supported is False
    assert result.media_kind == "text"
    assert result.synthid["status"] == "skipped"

    # 未知格式
    blob = tmp_path / "blob.xyz"
    blob.write_bytes(b"data")
    result = detect_deep_watermark(blob)
    assert result.supported is False
    assert result.media_kind == "unsupported"
    assert all("不伪装" in n for n in result.notes)


def test_deep_watermark_detect_synthid_can_be_enabled_via_env(tmp_path, monkeypatch):
    from PIL import Image

    from app.services.asset_provenance import detect_deep_watermark

    monkeypatch.setenv("YLCRAFT_SYNTHID_DETECT_ENABLED", "1")
    monkeypatch.setenv("YLCRAFT_SYNTHID_DETECT_PROVIDER", "demo-detector")
    source = tmp_path / "img.png"
    Image.new("RGB", (32, 32), "white").save(source)

    result = detect_deep_watermark(source, "image/png")
    assert result.synthid["status"] == "enabled"
    assert result.synthid["provider"] == "demo-detector"
    # 仍只读
    assert "未修改" in " ".join(result.notes)


def test_deep_watermark_detect_video_samples_frames_and_scores(tmp_path):
    """视频显性水印检测：抽取关键帧做统计检测，返回平均分，仍只读。"""
    import shutil
    import subprocess

    if not (shutil.which("ffmpeg") and shutil.which("ffprobe")):
        return  # 环境缺 ffmpeg 时跳过

    from app.services.asset_provenance import detect_deep_watermark

    video = tmp_path / "clip.mp4"
    subprocess.run(
        ["ffmpeg", "-y", "-hide_banner", "-loglevel", "error",
         "-f", "lavfi", "-i", "testsrc2=size=160x120:rate=5", "-t", "0.6", str(video)],
        check=True, capture_output=True,
    )
    before = video.read_bytes()
    result = detect_deep_watermark(video, "video/mp4")
    assert result.supported is True
    assert result.media_kind == "video"
    assert result.ctrlregen["method"].startswith("statistical-ctrlregen-like")
    assert result.ctrlregen["frame_count"] > 0
    assert video.read_bytes() == before  # 只读
    assert any("未修改" in n for n in result.notes)


def test_visual_watermark_remove_image_delogo_and_preserve_source(tmp_path):
    """图片 delogo：生成派生副本，源文件不被修改。"""
    import shutil
    import subprocess

    if not (shutil.which("ffmpeg") and shutil.which("ffprobe")):
        return

    from app.services.asset_provenance.visual_watermark import remove_visual_watermark_dict

    src = tmp_path / "wm.png"
    subprocess.run(
        ["ffmpeg", "-y", "-hide_banner", "-loglevel", "error",
         "-f", "lavfi", "-i", "color=c=blue:s=320x240:d=1", "-frames:v", "1", str(src)],
        check=True, capture_output=True,
    )
    before = src.read_bytes()
    result = remove_visual_watermark_dict(src, region={"corner": "top_right", "inset": 4}, method="delogo")
    assert result["supported"] is True
    assert result["media_kind"] == "image"
    assert result["method"] == "delogo"
    assert result["output_path"] and Path(result["output_path"]).exists()
    assert src.read_bytes() == before  # 源文件保留
    assert any("未被修改" in n for n in result["notes"])


def test_visual_watermark_remove_unsupported_format(tmp_path):
    """不支持格式：只报告，不伪装成已去除。"""
    from app.services.asset_provenance.visual_watermark import remove_visual_watermark_dict

    blob = tmp_path / "blob.xyz"
    blob.write_bytes(b"data")
    result = remove_visual_watermark_dict(blob, region={"corner": "top_right"}, method="delogo")
    assert result["supported"] is False
    assert result["media_kind"] == "unsupported"
    assert result["output_path"] == ""
